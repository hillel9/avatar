#!/usr/bin/env python3
"""Generate Avatar Videos Report PowerPoint from data.json."""

import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData

SCRIPT_DIR = Path(__file__).parent
DATA_FILE = SCRIPT_DIR / "data.json"
OUTPUT_FILE = SCRIPT_DIR / "avatar_videos_report.pptx"

# Colors
BLUE = RGBColor(0x15, 0x65, 0xC0)
DARK_BLUE = RGBColor(0x0D, 0x47, 0xA1)
PURPLE = RGBColor(0x7B, 0x1F, 0xA2)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xE6, 0x51, 0x00)
DARK = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x61, 0x61, 0x61)
LIGHT_GRAY = RGBColor(0x9E, 0x9E, 0x9E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)


def load_data():
    with open(DATA_FILE) as f:
        return json.load(f)


def set_slide_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 bold=False, color=DARK, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def make_title_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, WHITE)

    # Blue accent bar at top
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()

    add_text_box(slide, Inches(1), Inches(2.5), Inches(8), Inches(1),
                 "Avatar Videos Report", font_size=40, bold=True, color=DARK)
    add_text_box(slide, Inches(1), Inches(3.5), Inches(8), Inches(0.6),
                 f"Updated {data['updated_at']}", font_size=18, color=GRAY)
    add_text_box(slide, Inches(1), Inches(4.2), Inches(8), Inches(0.6),
                 "Snapshot of All Accounts", font_size=14, color=LIGHT_GRAY)


def make_kpi_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
                 "Snapshot — Last 30 Days", font_size=24, bold=True, color=DARK)

    # Date range from daily_users
    dates = data["daily_users"]
    first = dates[0]["date"].replace("2026-", "")
    last = dates[-1]["date"].replace("2026-", "")
    # Format dates
    from datetime import datetime
    d1 = datetime.strptime(dates[0]["date"], "%Y-%m-%d").strftime("%b %d")
    d2 = datetime.strptime(dates[-1]["date"], "%Y-%m-%d").strftime("%b %d")
    add_text_box(slide, Inches(0.8), Inches(0.95), Inches(6), Inches(0.4),
                 f"{d1} – {d2}", font_size=12, color=LIGHT_GRAY)

    kpis = [
        ("Monthly Unique Users", str(data["unique_users"]), BLUE),
        ("Monthly Avatar Videos", str(data["total_videos"]), PURPLE),
        ("Avg. Videos per User", str(data["avg_videos_per_user"]), GREEN),
        ("Returning Users", f"{data['returning_users_pct']}%", ORANGE),
    ]

    card_w = Inches(2.8)
    card_h = Inches(2.2)
    start_x = Inches(0.8)
    y = Inches(1.8)
    gap = Inches(0.25)

    for i, (label, value, color) in enumerate(kpis):
        x = start_x + i * (card_w + gap)
        # Card background
        card = slide.shapes.add_shape(1, x, y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.fill.background()
        card.shadow.inherit = False

        add_text_box(slide, x + Inches(0.3), y + Inches(0.4), Inches(2.2), Inches(0.4),
                     label, font_size=11, color=GRAY)
        add_text_box(slide, x + Inches(0.3), y + Inches(1.0), Inches(2.2), Inches(0.8),
                     value, font_size=36, bold=True, color=color)


def make_chart_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
                 "Unique Users Over Time", font_size=24, bold=True, color=DARK)

    chart_data = CategoryChartData()
    from datetime import datetime
    categories = []
    values = []
    for d in data["daily_users"]:
        dt = datetime.strptime(d["date"], "%Y-%m-%d")
        categories.append(dt.strftime("%b %d"))
        values.append(d["count"])

    chart_data.categories = categories
    chart_data.add_series("Users", values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.5),
        chart_data
    ).chart

    chart.has_legend = False
    plot = chart.plots[0]
    plot.smooth = True
    series = plot.series[0]
    series.format.line.color.rgb = BLUE
    series.format.line.width = Pt(2.5)
    series.smooth = True

    # Style axes
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(8)
    cat_axis.tick_labels.font.color.rgb = GRAY
    val_axis = chart.value_axis
    val_axis.tick_labels.font.size = Pt(9)
    val_axis.tick_labels.font.color.rgb = GRAY


def make_poc_table_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
                 "Customers POC — Account Activity", font_size=24, bold=True, color=DARK)
    add_text_box(slide, Inches(0.8), Inches(0.95), Inches(8), Inches(0.4),
                 "Accounts enabled in sandboxes, staging environments, or with limited creators.",
                 font_size=12, color=GRAY)

    # Aggregate POC accounts
    accounts = {}
    for acc in data["poc_accounts"]:
        name = acc["name"]
        # Normalize names
        if "AHK" in name:
            key = "AHK"
        elif "eLearning" in name:
            key = "eLearning Media"
        elif "EY" in name:
            key = "EY"
        elif "McDonald" in name:
            key = "McDonald's"
        elif "Michigan" in name:
            key = "U of Michigan"
        elif "Roche" in name:
            key = "Roche"
        elif "SAP" in name:
            key = "SAP"
        elif "Kaltura" in name:
            continue  # Skip internal
        else:
            key = name

        if key not in accounts:
            accounts[key] = {"users": 0, "videos": 0, "last": acc["last_active"]}
        accounts[key]["users"] += acc["active_users"]
        accounts[key]["videos"] += acc["videos_created"]
        if acc["last_active"] > accounts[key]["last"]:
            accounts[key]["last"] = acc["last_active"]

    # Get returning rates
    returning_map = {
        "eLearning Media": data["poc_returning"].get("eLearning Media", 0),
        "AHK": data["poc_returning"].get("AHK Sandbox", 0),
        "EY": data["poc_returning"].get("EY - Sandbox", 0),
        "McDonald's": data["poc_returning"].get("McDonald's Prod", 0),
        "Roche": data["poc_returning"].get("Roche Learning PROD", 0),
        "U of Michigan": data["poc_returning"].get("U of Michigan (KMC1)", 0),
        "SAP": data["poc_returning"].get("SAP Media Share", 0),
    }

    # Sort by users desc
    sorted_accounts = sorted(accounts.items(), key=lambda x: (-x[1]["users"], -x[1]["videos"]))

    headers = ["Customer", "Active Users", "Returning", "Videos", "Last Login"]
    rows = len(sorted_accounts) + 1
    cols = len(headers)

    table_shape = slide.shapes.add_table(rows, cols, Inches(0.6), Inches(1.5), Inches(11.5), Inches(0.4 * rows))
    table = table_shape.table

    # Set column widths
    table.columns[0].width = Inches(3.0)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(2.0)
    table.columns[3].width = Inches(2.0)
    table.columns[4].width = Inches(2.5)

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE

    # Data rows
    from datetime import datetime
    for row_idx, (name, info) in enumerate(sorted_accounts, 1):
        ret = returning_map.get(name, 0)
        last_dt = datetime.strptime(info["last"], "%Y-%m-%d").strftime("%b %d")
        values = [name, str(info["users"]), f"{ret}%", str(info["videos"]), last_dt]
        for col_idx, val in enumerate(values):
            cell = table.cell(row_idx, col_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = DARK
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)


def make_free_trial_slide(prs, data, cohort_num, cohort_accounts, launch_label):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                 f"Free Trial Program — Cohort {cohort_num}", font_size=24, bold=True, color=DARK)
    add_text_box(slide, Inches(0.8), Inches(0.95), Inches(8), Inches(0.4),
                 launch_label, font_size=12, color=LIGHT_GRAY)

    # Build lookup from free_trial data
    trial_lookup = {}
    for t in data["free_trial"]:
        trial_lookup[t["name"]] = t

    headers = ["Customer", "Saw Modal", "Entered App", "Returning", "Videos", "Potential"]
    rows = len(cohort_accounts) + 1
    cols = len(headers)

    table_shape = slide.shapes.add_table(rows, cols, Inches(0.4), Inches(1.5), Inches(12.2), Inches(0.4 * rows))
    table = table_shape.table

    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(1.8)
    table.columns[2].width = Inches(1.8)
    table.columns[3].width = Inches(1.8)
    table.columns[4].width = Inches(1.5)
    table.columns[5].width = Inches(1.8)

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE

    for row_idx, (display_name, data_key) in enumerate(cohort_accounts, 1):
        modal = data["free_trial_modal"].get(data_key, 0)
        trial = trial_lookup.get(data_key, {})
        entered = trial.get("active_users", 0)
        returning = data["free_trial_returning"].get(data_key, 0)
        videos = trial.get("videos_created", 0)
        potential = data["free_trial_potential"].get(data_key, 0)

        values = [
            display_name,
            f"{modal:,}",
            str(entered),
            f"{returning}%",
            str(videos),
            f"{potential:,}",
        ]
        for col_idx, val in enumerate(values):
            cell = table.cell(row_idx, col_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = DARK
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)


def make_methods_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
                 "Video Creation Methods", font_size=24, bold=True, color=DARK)
    add_text_box(slide, Inches(0.8), Inches(0.95), Inches(8), Inches(0.4),
                 "Data from POC customers and Free Trial accounts only.",
                 font_size=12, color=GRAY)

    methods = data["video_methods"]
    chart_data = CategoryChartData()
    # Sort by value desc
    sorted_methods = sorted(methods.items(), key=lambda x: -x[1])
    categories = [m[0] for m in sorted_methods]
    values = [m[1] for m in sorted_methods]

    chart_data.categories = categories
    chart_data.add_series("Videos", values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.8), Inches(1.5), Inches(11), Inches(5),
        chart_data
    ).chart

    chart.has_legend = False
    plot = chart.plots[0]
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = BLUE

    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.color.rgb = DARK
    val_axis = chart.value_axis
    val_axis.tick_labels.font.size = Pt(10)
    val_axis.tick_labels.font.color.rgb = GRAY


def make_content_analysis_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                 "Free Trial Content Analysis", font_size=24, bold=True, color=DARK)
    add_text_box(slide, Inches(0.8), Inches(0.95), Inches(10), Inches(0.4),
                 "Sample of video scripts generated by Free Trial users.",
                 font_size=12, color=GRAY)

    scripts = [
        ("Contract Signing Authority", "Compliance training from source recording", "~3-4 min", "Compliance Training"),
        ("Dr. Dickens Welcome", "Instructor intro from webcam recording", "~2-3 min", "Course Welcome"),
        ("Dr. Gregory Welcome", "Motivational course intro from recording", "~3-4 min", "Course Welcome"),
        ("Supplements", "Student assignment, slide-based research", "~4-5 min", "Student Assignment"),
        ("Healthcare System", "Student academic presentation", "~5-6 min", "Student Assignment"),
        ("BuyLU Guide", "Staff portal walkthrough", "~30-45 sec", "Process Training"),
        ("Prompt Together", "AI prompting lesson segment", "~1 min", "Course Content"),
        ("Universe Intro", "Personal class greeting", "~20-30 sec", "Course Welcome"),
        ("NEPA", "Abandoned student assignment", "~10 sec", "Abandoned"),
        ("System Test", "Pipeline verification", "~5 sec", "Test"),
    ]

    headers = ["Video", "Intent", "Length", "Category"]
    rows = len(scripts) + 1
    cols = len(headers)

    table_shape = slide.shapes.add_table(rows, cols, Inches(0.4), Inches(1.4), Inches(12.3), Inches(0.37 * rows))
    table = table_shape.table

    table.columns[0].width = Inches(2.8)
    table.columns[1].width = Inches(4.5)
    table.columns[2].width = Inches(2.0)
    table.columns[3].width = Inches(3.0)

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE

    for row_idx, (video, intent, length, category) in enumerate(scripts, 1):
        values = [video, intent, length, category]
        for col_idx, val in enumerate(values):
            cell = table.cell(row_idx, col_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = DARK
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)


def main():
    data = load_data()

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    make_title_slide(prs, data)
    make_kpi_slide(prs, data)
    make_chart_slide(prs, data)
    make_poc_table_slide(prs, data)

    # Cohort 1
    cohort1 = [
        ("Liberty University", "Liberty University"),
        ("University of Maine System", "University of Maine System"),
        ("Block, Inc.", "Block , Inc."),
        ("Allspring Global Investments", "Allspring Global Investments"),
    ]
    make_free_trial_slide(prs, data, 1, cohort1, "Launched Jun 8–12")

    # Cohort 2
    cohort2 = [
        ("Colorado Christian University", "Colorado Christian University"),
        ("University of Arkansas Fayetteville", "UARK"),
        ("Saskatchewan Polytechnic", "2167551 - Sask Polytech - PROD"),
        ("California State University, Long Beach", "California State University, Long Beach"),
        ("Olivet Nazarene University", "Olivet Nazarene University - Production"),
    ]
    make_free_trial_slide(prs, data, 2, cohort2, "Launched Jun 22–26")

    make_methods_slide(prs, data)
    make_content_analysis_slide(prs)

    prs.save(str(OUTPUT_FILE))
    print(f"Saved: {OUTPUT_FILE}")


if __name__ == "__main__":
    main()
