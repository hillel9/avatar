#!/usr/bin/env python3
"""Generate a single-slide PowerPoint with Use Case Breakdown."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

OUTPUT_FILE = Path(__file__).parent / "use_case_breakdown.pptx"

# Colors
BLUE = RGBColor(0x15, 0x65, 0xC0)
PURPLE = RGBColor(0x7B, 0x1F, 0xA2)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
GRAY_BAR = RGBColor(0x9E, 0x9E, 0x9E)
DARK = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x61, 0x61, 0x61)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Use Case Breakdown"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK
    p.font.name = "Calibri"

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(10), Inches(0.4))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Categorization of videos created by Free Trial users (sample of published content)"
    p2.font.size = Pt(12)
    p2.font.color.rgb = GRAY
    p2.font.name = "Calibri"

    # Use case data
    use_cases = [
        ("Course Welcome / Intro", 4, BLUE),
        ("Student Assignments", 3, PURPLE),
        ("Training / Process", 2, GREEN),
        ("Test / Abandoned", 1, GRAY_BAR),
    ]

    # Create horizontal bar chart
    chart_data = CategoryChartData()
    categories = [uc[0] for uc in use_cases]
    values = [uc[1] for uc in use_cases]

    chart_data.categories = categories
    chart_data.add_series("Videos", values)

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.8), Inches(1.8), Inches(11), Inches(5.2),
        chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = False

    plot = chart.plots[0]
    plot.gap_width = 100

    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = BLUE

    # Color each point individually
    colors = [BLUE, PURPLE, GREEN, GRAY_BAR]
    for i, color in enumerate(colors):
        pt = series.points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = color

    # Style axes
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(14)
    cat_axis.tick_labels.font.color.rgb = DARK
    cat_axis.tick_labels.font.name = "Calibri"
    cat_axis.has_major_gridlines = False

    val_axis = chart.value_axis
    val_axis.tick_labels.font.size = Pt(11)
    val_axis.tick_labels.font.color.rgb = GRAY
    val_axis.has_major_gridlines = True
    val_axis.major_gridlines.format.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    val_axis.maximum_scale = 5

    # Add data labels
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.show_value = True
    data_labels.font.size = Pt(14)
    data_labels.font.bold = True
    data_labels.font.color.rgb = DARK

    prs.save(str(OUTPUT_FILE))
    print(f"Saved: {OUTPUT_FILE}")


if __name__ == "__main__":
    main()
